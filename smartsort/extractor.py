"""Text and metadata extraction for supported document types."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class ExtractionResult:
    filename: str
    text: str
    metadata: dict[str, str]


class DocumentExtractor:
    """Extract a bounded amount of first-page/first-slide text and metadata."""

    max_characters = 12_000

    def extract(self, path: Path) -> ExtractionResult:
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            text, metadata = self._pdf(path)
        elif suffix == ".docx":
            text, metadata = self._docx(path)
        elif suffix == ".pptx":
            text, metadata = self._pptx(path)
        elif suffix == ".txt":
            text, metadata = self._txt(path)
        else:
            raise ValueError(f"Unsupported file type: {path.suffix}")
        return ExtractionResult(path.name, text[: self.max_characters], metadata)

    @staticmethod
    def _pdf(path: Path) -> tuple[str, dict[str, str]]:
        import fitz
        with fitz.open(path) as document:
            text = document[0].get_text("text") if document.page_count else ""
            metadata = {str(k): str(v) for k, v in document.metadata.items() if v}
        return text, metadata

    @staticmethod
    def _docx(path: Path) -> tuple[str, dict[str, str]]:
        from docx import Document
        document = Document(path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        props = document.core_properties
        metadata = {key: str(getattr(props, key)) for key in ("title", "author", "subject", "keywords") if getattr(props, key, None)}
        return text, metadata

    @staticmethod
    def _pptx(path: Path) -> tuple[str, dict[str, str]]:
        from pptx import Presentation
        presentation = Presentation(path)
        first_slide = presentation.slides[0] if presentation.slides else []
        text = "\n".join(shape.text for shape in first_slide if hasattr(shape, "text"))
        props = presentation.core_properties
        metadata = {key: str(getattr(props, key)) for key in ("title", "author", "subject", "keywords") if getattr(props, key, None)}
        return text, metadata

    @staticmethod
    def _txt(path: Path) -> tuple[str, dict[str, str]]:
        try:
            return path.read_text(encoding="utf-8")[:12_000], {}
        except UnicodeDecodeError:
            return path.read_text(encoding="latin-1")[:12_000], {}
